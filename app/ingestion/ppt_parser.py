from pptx import Presentation


class PPTParser:

    def __init__(self, ppt_path):
        self.ppt_path = ppt_path

    def extract_slides(self):

        prs = Presentation(self.ppt_path)

        slides_data = []

        for idx, slide in enumerate(prs.slides):

            slide_text = []

            for shape in slide.shapes:

                if hasattr(shape, "text"):
                    slide_text.append(shape.text)

            slides_data.append(
                {
                    "slide_number": idx + 1,
                    "content": "\n".join(slide_text)
                }
            )

        return slides_data